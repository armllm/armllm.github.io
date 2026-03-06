from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUTPUT_FILE = "2026/Armenia-LLM-Summer-School-2026.pptx"


BG_COLOR = RGBColor(11, 18, 32)
TITLE_COLOR = RGBColor(129, 140, 248)
TEXT_COLOR = RGBColor(240, 244, 255)
MUTED_COLOR = RGBColor(193, 201, 219)
ACCENT_COLOR = RGBColor(79, 70, 229)


def set_dark_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.55))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT_COLOR
    band.line.fill.background()


def add_slide_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.75), Inches(12), Inches(1.2))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.7), Inches(12), Inches(0.8))
        sub_tf = sub_box.text_frame
        sub_tf.clear()
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle
        sub_p.font.size = Pt(18)
        sub_p.font.color.rgb = MUTED_COLOR


def add_bullets(slide, bullets, left=0.9, top=2.4, width=11.8, height=4.6, font_size=24):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(10)


def add_two_column_bullets(slide, left_title, left_bullets, right_title, right_bullets):
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.15), Inches(6.15), Inches(4.7))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(20, 28, 49)
    left_box.line.fill.background()

    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.95), Inches(2.15), Inches(5.7), Inches(4.7))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(20, 28, 49)
    right_box.line.fill.background()

    left_title_shape = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(5.5), Inches(0.5))
    left_title_tf = left_title_shape.text_frame
    left_title_tf.text = left_title
    left_title_tf.paragraphs[0].font.size = Pt(20)
    left_title_tf.paragraphs[0].font.bold = True
    left_title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    right_title_shape = slide.shapes.add_textbox(Inches(7.2), Inches(2.35), Inches(5.2), Inches(0.5))
    right_title_tf = right_title_shape.text_frame
    right_title_tf.text = right_title
    right_title_tf.paragraphs[0].font.size = Pt(20)
    right_title_tf.paragraphs[0].font.bold = True
    right_title_tf.paragraphs[0].font.color.rgb = TITLE_COLOR

    add_bullets(slide, left_bullets, left=1.0, top=2.9, width=5.6, height=3.7, font_size=18)
    add_bullets(slide, right_bullets, left=7.2, top=2.9, width=5.0, height=3.7, font_size=18)


def add_footer(slide, text):
    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.35))
    tf = footer.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED_COLOR


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(
        slide,
        "Armenia LLM Summer School 2026",
        "Based on https://armllm.github.io/2026/",
    )
    add_bullets(
        slide,
        [
            "AI9 Startup Campus • Yerevan, Armenia",
            "August 3–7, 2026 (Hackathon: August 8, 24-hour)",
            "In-person program + Zoom-based online theory track",
            "Applications are currently open",
        ],
        top=2.35,
        font_size=25,
    )
    add_footer(slide, "Armenia LLM Summer School • 2026 Overview")

    # Slide 2
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Program Overview", "Annual event focused on modern LLM education")
    add_bullets(
        slide,
        [
            "Comprehensive lectures on Large Language Models",
            "Hands-on workshops and practical sessions",
            "Online theoretical track for broader access",
            "24-hour AI9-organized hackathon with ArmLLM co-organizing content",
            "Built on successful 2024 and 2025 editions",
        ],
        top=2.3,
        font_size=23,
    )
    add_footer(slide, "Source: About + Format sections on armllm.github.io/2026/")

    # Slide 3
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Key Dates and Location")
    add_two_column_bullets(
        slide,
        "Important Dates",
        [
            "May 15, 2026: Application deadline (form + CV)",
            "May 28, 2026: Technical assessment sent by email",
            "By June 5, 2026: Admission decisions",
            "Aug 3–7, 2026: Summer school week",
            "Aug 8, 2026: Hackathon (starts 00:00)",
        ],
        "Venue & Format",
        [
            "Venue: AI9 Startup Campus",
            "Address: 9 Isahakyan St, Yerevan, Armenia",
            "Hybrid attendance model",
            "In-person includes practical sessions and GPU access",
            "Online focuses on theory via Zoom",
        ],
    )
    add_footer(slide, "Source: Hero, About, Apply, and FAQ sections")

    # Slide 4
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Application Process")
    add_bullets(
        slide,
        [
            "Two-stage competitive selection process",
            "Stage 1: Submit application form + CV",
            "Stage 2: Complete technical assessment",
            "Watch email for test link and official updates",
            "Contact: armeniallm@gmail.com",
        ],
        top=2.35,
        font_size=24,
    )
    add_footer(slide, "Apply links: forms.gle/1HnZ4BLitQkXo2reA and fee waiver form")

    # Slide 5
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Fees and Financial Aid")
    add_two_column_bullets(
        slide,
        "Participation Fees",
        [
            "In-person: 100,000 AMD (or equivalent)",
            "Online: $200 (or equivalent)",
            "In-person includes theory + practical + GPU access",
            "Online typically does not include shared GPU access",
        ],
        "Financial Aid",
        [
            "Student fee waivers up to 90%",
            "Additional discounts may be available based on funding",
            "Aid request can be submitted in the application flow",
            "International payment support provided after acceptance",
        ],
    )
    add_footer(slide, "Source: Fees section + FAQ (fees, aid, payments)")

    # Slide 6
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Participant Experience")
    add_two_column_bullets(
        slide,
        "What Participants Need",
        [
            "Bring a laptop for all practical sessions",
            "Python proficiency is expected",
            "Foundational ML / neural network knowledge is expected",
            "No prior hands-on LLM experience is required",
        ],
        "What Participants Receive",
        [
            "Full-week learning across theory and practice",
            "Slack channel for daily agenda and materials",
            "Certificate of Participation (upon full attendance)",
            "Hackathon participation pathway (separate registration)",
        ],
    )
    add_footer(slide, "Source: FAQ (overview, eligibility, participation)")

    # Slide 7
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Speakers, Agenda, and Current Status")
    add_bullets(
        slide,
        [
            "Featured speakers list: To be updated",
            "Day-by-day agenda: To be updated",
            "Applications and logistics information are already published",
            "Program details will be announced closer to event dates",
        ],
        top=2.5,
        font_size=24,
    )
    add_footer(slide, "Source: Speakers and Agenda sections on the 2026 page")

    # Slide 8
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Organizers and Sponsor")
    add_two_column_bullets(
        slide,
        "Organizers",
        [
            "Armen Aghajanyan (Perceptron AI)",
            "Erik Arakelyan (Nvidia)",
            "Aram Galstyan (USC / Amazon)",
            "Hrant Khachatryan (YerevaNN)",
            "Perouz Taslakian (ServiceNow Research)",
            "Tatevik Vardanyan (YerevaNN)",
        ],
        "2026 Sponsor",
        [
            "Venue Sponsor: AI9 Startup Campus",
            "Partner website: https://www.ai9.am/",
            "Location support for in-person activities",
            "Hackathon organized by AI9",
        ],
    )
    add_footer(slide, "Source: Sponsors + Organizers sections")

    # Slide 9
    slide = prs.slides.add_slide(blank)
    set_dark_background(slide)
    add_slide_title(slide, "Past Editions and Links")
    add_bullets(
        slide,
        [
            "2025 (2nd edition): Jul 24–30, 2025",
            "2024 (1st edition): Inaugural summer school",
            "2026 page: https://armllm.github.io/2026/",
            "Apply: https://forms.gle/1HnZ4BLitQkXo2reA",
            "Fee waiver: https://forms.gle/NMMsdenubnoGkRJ87",
            "Contact: armeniallm@gmail.com",
        ],
        top=2.35,
        font_size=23,
    )
    add_footer(slide, "Armenia LLM Summer School • Deck generated from website content")

    prs.save(OUTPUT_FILE)


if __name__ == "__main__":
    main()
